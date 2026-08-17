"""PDF do deck -> PowerPoint, uma pagina por slide, sem perder nada do design.

    python ferramentas/pdf_para_pptx.py docs/Deck-Smart-Store.pdf

POR QUE PASSA PELO PDF.

Converter HTML direto em PowerPoint exige um navegador para rasterizar as
paginas — e um navegador e a unica coisa capaz de honrar o CSS moderno do deck
(oklch, grid, flex, fontes web). Reconstruir os slides "parecidos" em PowerPoint
nao e conversao: e outro deck.

    Preservar o design significa preservar os PIXELS. Tudo o que redesenha,
    muda.

Entao o navegador faz a parte que so ele sabe fazer (Ctrl+P -> salvar como PDF),
e este script faz o resto: cada pagina vira uma imagem em resolucao de
apresentacao e ocupa um slide inteiro, com as notas do apresentador no lugar
certo.
"""
import subprocess
import sys
import tempfile
from pathlib import Path


def notas_do_html(html):
    """As falas ficam em data-speaker-notes no HTML, na ordem das secoes."""
    import html as H
    import re
    if not html or not Path(html).exists():
        return []
    t = Path(html).read_text(encoding="utf-8")
    return [H.unescape(x) for x in re.findall(r'data-speaker-notes="([^"]*)"', t)]


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        return
    pdf = Path(sys.argv[1]).resolve()
    if not pdf.exists():
        print(f"nao achei {pdf}")
        return
    html = pdf.with_suffix(".html")
    saida = pdf.with_suffix(".pptx")

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print("falta a biblioteca:  pip install python-pptx")
        return

    with tempfile.TemporaryDirectory() as tmp:
        # 150 dpi: nitido em projetor e ainda leve o bastante para enviar.
        subprocess.run(["pdftoppm", "-png", "-r", "150", str(pdf),
                        str(Path(tmp) / "p")], check=True)
        paginas = sorted(Path(tmp).glob("p-*.png"))
        if not paginas:
            print("pdftoppm nao gerou nada — o Poppler esta instalado?")
            return

        falas = notas_do_html(html)
        pr = Presentation()
        # 16:9 em EMU. O slide inteiro e a imagem: sem margem, sem sobra.
        pr.slide_width, pr.slide_height = Emu(12192000), Emu(6858000)
        vazio = pr.slide_layouts[6]

        for i, img in enumerate(paginas):
            s = pr.slides.add_slide(vazio)
            s.shapes.add_picture(str(img), 0, 0,
                                 width=pr.slide_width, height=pr.slide_height)
            if i < len(falas) and falas[i]:
                s.notes_slide.notes_text_frame.text = falas[i]

        pr.save(str(saida))

    print(f"  {len(paginas)} slides  ->  {saida}")
    if falas:
        print(f"  {min(len(falas), len(paginas))} notas do apresentador incluidas")
if __name__ == "__main__":
    main()
